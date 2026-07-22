"""
master_export.py

PowerPoint deck builder for the "master module" -- takes already-fetched
ACS/BLS dashboard dicts (same {chart_type, ...} shape demographics_dashboard.py/
bls_dashboard.py already produce) plus optional parsed CoStar/SmartRE data,
and returns a python-pptx Presentation with one native, editable chart per
slide. Same "server-side native chart generation" architecture as the
Excel export (dashboard_excel_export.py / excel_export.py), just targeting
python-pptx instead of openpyxl -- confirmed via a real save/reopen/
read-back round trip that python-pptx's add_chart() embeds a genuine OOXML
chart object, not a rasterized image.

Chart-type dispatch mirrors dashboard_excel_export.py's _add_block()
exactly: "line"/"multi_line" -> LINE_MARKERS, "stacked_bar" -> COLUMN_STACKED,
"bar" -> COLUMN_CLUSTERED. Unlike openpyxl (which needs a worksheet + cell
Reference for chart data), python-pptx's CategoryChartData carries category/
series values directly -- no intermediate "write a table" step needed.

Every deck is built from assets/master_deck_template.pptx -- a ~1MB derived
copy of the user's real reference deck ("Report Builder/Reference - Sandy
Springs Comp Plan...pptx", 39MB with 82 hand-built slides) with every
content slide and its now-orphaned media/charts/embeddings stripped out
(via the pptx skill's clean.py), keeping only the theme/masters/layouts --
confirmed via validate.py against the original that nothing load-bearing
was lost. Loading the full 39MB reference deck per web request was ruled
out (tested: removing all its slides with python-pptx alone does NOT
shrink the file -- orphaned parts aren't auto-pruned on save), so the
lean template is the actual runtime dependency; the two source files stay
in "Report Builder/" for reference only, not read at runtime.

Layout choices, confirmed by inspecting how the reference deck's own real
slides actually used them: "Title Slide" for the cover, "3_Large Statement
Teal" for section dividers (this is the exact layout the reference deck's
own "Executive Summary"/"Demographic Trends"/etc. dividers used), "Title
Only" for chart slides (a clean single Title placeholder to hold the
section name as a heading, plus that layout's body placeholder used here
as a subheading for the chart's own metric name).

Color scheme: the layouts above all live on the reference deck's own
master/theme (ppt/theme/theme1.xml), but that particular theme part's own
color scheme was just generic unbranded Office defaults (accent1 a plain
blue) -- confirmed the real brand colors live in a separate file,
"Report Builder/KB_Standard_PPTX_Theme.thmx" (its own theme1.xml, color
scheme named "HousingTrustFund", accent1 a dark forest green). The
template's theme1.xml has that <a:clrScheme> block swapped in for
KB_Standard's, keeping every layout/shape/font untouched -- only the
color values changed.

Every chart slide's chart sits flush right (per explicit feedback) --
earlier versions alternated left/right per slide; that alternation is
gone, every chart now uses the same right-hand geometry.

Known simplification (v1): python-pptx exposes no display-blanks-as toggle
(unlike openpyxl's chart.display_blanks = "span", added for the Excel
export after a real gap-rendering bug). A None value in a series renders as
a gap here rather than a connected line across it. ACS/BLS data rarely has
mid-range gaps in practice; revisit if real decks surface this.

Known simplification (v2, SmartRE): the Excel export's full 12-cut
breakdown (Overall/New/Resale x Single-Family/Townhome/Condo) and its
sale-price scatter chart are both deliberately not reproduced here -- a
report deck's audience wants one headline picture, not the full drill-down
(still available via the dedicated SmartRE Excel export), and python-pptx's
XyChartData has no equivalent to openpyxl's date-typed category axis (it
only accepts plain numeric x-values), so a real calendar-date x-axis isn't
achievable without writing raw date serials as the axis labels. This
module includes only the "Overall" cut's year x price-bin stacked bar.

Known simplification (v3, CoStar Market Overview class routing): every
uploaded property class (including Multifamily) lands in Commercial Real
Estate Analysis here, unlike the plan's original Housing-Analysis-for-
Multifamily split -- keeping all Market Overview output on one code path
was judged not worth fragmenting for a single class; revisit if this reads
oddly in a real deck.
"""

from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from bls_dashboard import NAICS_SECTORS
from costar_heartbeat import BROAD_CLASSES, _decade_class_table
from costar_market_overview import CLASS_LABELS, METRICS_BY_CLASS, PROPERTY_CLASSES, _is_rate_metric, parse_grid
from dashboard_excel_export import acs_chart_title, bls_chart_title
from smartre_sales import PRICE_BIN_LABELS, price_bin_label

TEMPLATE_PATH = "assets/master_deck_template.pptx"

# ACS chart keys per named section (fixed -- ACS's chart set doesn't
# change per request, unlike BLS's per-sector-selectable keys below).
# "All BLS data" (per the source doc) always lands in Economic Analysis
# unconditionally, so BLS needs no separate per-section key list here.
ACS_SECTION_CHARTS = {
    "Demographic Analysis": [
        "population", "households", "age_by_cohort", "age_by_cohort_simplified",
        "median_age", "race", "hispanic_ethnicity", "household_size", "household_type",
    ],
    "Economic Analysis": [
        "household_income", "household_income_simplified", "median_household_income",
    ],
    "Housing Analysis": [
        "housing_units", "housing_unit_occupancy", "housing_unit_type", "housing_unit_type_simplified",
        "year_built", "year_moved_in", "tenure", "median_home_value", "median_rent",
        "owner_cost_burden", "renter_cost_burden", "tenure_by_age_owner", "tenure_by_age_renter",
        "tenure_by_income_owner", "tenure_by_income_renter",
    ],
}

_CHART_TYPE_MAP = {
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "multi_line": XL_CHART_TYPE.LINE_MARKERS,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
}

# Value-axis number formats, applied per chart rather than left raw --
# unlike the Excel export (dashboard_excel_export.py), which deliberately
# ships unformatted numbers for the user to format in a working
# spreadsheet, a presentation deck's axis should already read cleanly.
# "percent_fraction" is a real Office percentage format (multiplies the
# underlying 0-1 value by 100 for display) -- used for ACS stacked_bar/bar
# charts, whose `categories` values are genuine 0-1 fractions (confirmed in
# frontend/src/lib/chartMeta.ts's comment: "All stacked_bar charts... are
# percentages (0-1 fractions)"). "percent_scaled" is a custom format that
# only appends a literal "%" with no multiplication -- used for CoStar's
# own rate metrics (Vacancy/Occupancy), whose raw values are already on a
# 0-100 scale (confirmed in costar_market_overview.py's _to_number(),
# which strips a literal "%" off strings like "5.2%" to get 5.2, not
# 0.052) -- applying "percent_fraction" there would 100x-inflate them.
_AXIS_NUMBER_FORMATS = {
    "count": "#,##0",
    "dollars": "$#,##0",
    "years": "0.0",
    "percent_fraction": "0%",
    "percent_scaled": '0.0"%"',
}

# Mirrors frontend/src/lib/chartMeta.ts's CHART_META format field (hand-
# synced across the JS/Python boundary, same convention already used for
# ACS_CHART_TITLES above and MasterModule.tsx's ACS_SECTIONS/BLS_SECTIONS).
ACS_CHART_FORMAT = {
    "population": "count", "households": "count", "housing_units": "count",
    "housing_unit_occupancy": "percent_fraction", "housing_unit_type": "percent_fraction",
    "housing_unit_type_simplified": "percent_fraction", "year_built": "percent_fraction",
    "year_moved_in": "percent_fraction", "tenure": "percent_fraction",
    "median_home_value": "dollars", "median_rent": "dollars",
    "owner_cost_burden": "percent_fraction", "renter_cost_burden": "percent_fraction",
    "age_by_cohort": "percent_fraction", "age_by_cohort_simplified": "percent_fraction",
    "median_age": "years", "race": "percent_fraction", "hispanic_ethnicity": "percent_fraction",
    "household_income": "percent_fraction", "household_income_simplified": "percent_fraction",
    "median_household_income": "dollars",
    "tenure_by_age_owner": "percent_fraction", "tenure_by_age_renter": "percent_fraction",
    "tenure_by_income_owner": "percent_fraction", "tenure_by_income_renter": "percent_fraction",
    "household_size": "percent_fraction", "household_type": "percent_fraction",
}


def _bls_chart_format(key: str) -> str:
    """Mirrors frontend/src/lib/blsChartMeta.ts's blsChartMeta() -- BLS
    values are always a headcount or a dollar figure, never a percent."""
    if key in ("avg_pay_by_sector", "total_avg_pay_trend"):
        return "dollars"
    if key.startswith("wage_trend_") or key.startswith("avg_pay_trend_"):
        return "dollars"
    return "count"


def _market_metric_format(label: str) -> str:
    """CoStar Market Overview metric labels aren't a fixed key set (they're
    the sub-table headers from costar_market_overview.METRICS_BY_CLASS), so
    format is derived from the label text the same way that module's own
    _is_rate_metric() already derives chart shape from it."""
    lower = label.lower()
    if "vacan" in lower or "occupancy" in lower:
        return "percent_scaled"
    if "rent" in lower or lower in ("adr", "revpar"):
        return "dollars"
    return "count"

# Chart geometry (13.333" x 7.5" template): every chart sits flush right,
# a 6.0"-wide chart with a 0.5" right margin -- 13.333 - 0.5 - 6.0 = 6.833".
_CHART_TOP = Inches(1.7)
_CHART_HEIGHT = Inches(5.3)
_CHART_WIDTH = Inches(6.0)
_CHART_X = Inches(6.833)

# Confirmed by inspecting how the reference deck's own real section-divider
# slides positioned this placeholder (the layout's own default position is
# a leftover off-slide draft position, not what any real slide used).
_DIVIDER_TITLE_BOX = (Inches(0.396), Inches(1.362), Inches(12.767), Inches(2.269))
_DIVIDER_FONT_SIZE = Pt(44)

# Title slide geometry -- per explicit feedback, ported directly from
# "Report Builder/KB_ALT_PPTX_Theme.thmx"'s own "Section Header" layout
# (its title placeholder's exact off/ext box and 60pt font size), not
# derived from the divider box like the title slide's font size used to be.
# The title slide itself uses this deck's own "Title Only" layout (plain
# white background, no wave art) -- also per explicit feedback, picked as
# closest of the KB-branded candidates reviewed.
_TITLE_BOX = (Emu(831850), Emu(1709738), Emu(10515600), Emu(2852737))
_TITLE_FONT_SIZE = Pt(60)
_TITLE_BAR_HEIGHT = Inches(0.12)
_TITLE_BAR_GAP = Inches(0.15)


def _layout(prs: Presentation, name: str):
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"No slide layout named {name!r} in {TEMPLATE_PATH}")


def _placeholder(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _add_chart_slide(prs: Presentation, section: str, metric_title: str, chart_type: str, categories: list, series: dict, value_format: str = "count"):
    """One title + one native chart per slide (matches the "real, editable"
    requirement better than several charts per slide), chart flush right so
    the slide reads as a real report page, not a chart floating on blank
    space. section becomes the slide's heading (so every slide names which
    analysis it belongs to); metric_title becomes a subheading -- the chart
    itself carries no redundant internal title. value_format: a key into
    _AXIS_NUMBER_FORMATS, applied to the value axis so counts/dollars/
    percents read cleanly instead of as raw floats."""
    layout = _layout(prs, "Title Only")
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text_frame.text = section
    subtitle = _placeholder(slide, 15)
    if subtitle is not None:
        subtitle.text_frame.text = metric_title

    data = CategoryChartData()
    data.categories = [str(c) for c in categories]
    for name, values in series.items():
        data.add_series(name, values)

    xl_type = _CHART_TYPE_MAP[chart_type]
    graphic_frame = slide.shapes.add_chart(xl_type, _CHART_X, _CHART_TOP, _CHART_WIDTH, _CHART_HEIGHT, data)

    chart = graphic_frame.chart
    chart.has_title = False
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        # Charts only get a half-slide's width now (per the requested
        # layout), so python-pptx's default legend placement (overlapping
        # the plot area) reads far worse than it did at near-full-slide
        # width -- bottom + include_in_layout=False reserves real space
        # for it instead of layering it over the bars/lines.
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.value_axis.tick_labels.number_format = _AXIS_NUMBER_FORMATS.get(value_format, "#,##0")
    chart.value_axis.tick_labels.number_format_is_linked = False
    return slide


def _line_chart_slide(prs: Presentation, section: str, title: str, series: dict, value_format: str = "count") -> None:
    """series: {year: value} -- a single-series 'line' chart (ACS's
    population/median-* charts, BLS's per-sector/total trend charts)."""
    years = sorted(series)
    _add_chart_slide(prs, section, title, "line", years, {title: [series.get(y) for y in years]}, value_format)


def _multi_line_chart_slide(prs: Presentation, section: str, title: str, series_by_label: dict, value_format: str = "count") -> None:
    """series_by_label: {label: {year: value}} -- BLS's employment_by_sector/
    avg_pay_by_sector combined multi-series charts."""
    years = sorted({y for s in series_by_label.values() for y in s})
    series = {label: [series_by_label[label].get(y) for y in years] for label in series_by_label}
    _add_chart_slide(prs, section, title, "multi_line", years, series, value_format)


def _category_chart_slide(prs: Presentation, section: str, title: str, categories: dict, kind: str, value_format: str = "count") -> None:
    """categories: {year: {label: value}} -- ACS's stacked_bar/bar charts
    (category_breakdown()-shaped results)."""
    years = sorted(categories)
    labels = list(dict.fromkeys(label for y in years for label in categories[y]))
    series = {label: [categories[y].get(label) for y in years] for label in labels}
    _add_chart_slide(prs, section, title, kind, years, series, value_format)


def add_dashboard_chart_slide(prs: Presentation, section: str, title: str, chart: dict, value_format: str = "count") -> None:
    """Dispatches one ACS/BLS chart dict (the same {chart_type, ...} shape
    demographics_dashboard.py/bls_dashboard.py already return) to the right
    slide builder -- the single entry point build_master_deck's section
    loop calls per selected chart key."""
    chart_type = chart["chart_type"]
    if chart_type == "line":
        _line_chart_slide(prs, section, title, chart["series"], value_format)
    elif chart_type == "multi_line":
        _multi_line_chart_slide(prs, section, title, chart["series_by_label"], value_format)
    elif chart_type in ("stacked_bar", "bar"):
        _category_chart_slide(prs, section, title, chart["categories"], chart_type, value_format)


def _add_title_slide(prs: Presentation, title_text: str) -> None:
    """"Title Only" (this deck's own plain, no-wave-art layout -- also used
    for every chart slide) rather than the "Title Slide" layout's wave
    background, per explicit feedback picking it as the closest match to a
    reviewed KB-branded candidate. Title box position/size (60pt, left-
    aligned) is ported from KB_ALT_PPTX_Theme.thmx's own "Section Header"
    layout -- see _TITLE_BOX's own comment. A solid accent-color bar sits
    just under the title text, per explicit feedback ("a green bar across
    the page under the title slide text")."""
    layout = _layout(prs, "Title Only")
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    left, top, width, height = _TITLE_BOX
    title.left, title.top, title.width, title.height = left, top, width, height
    title.text_frame.word_wrap = True
    title.text_frame.text = title_text
    paragraph = title.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.bold = True
    run.font.size = _TITLE_FONT_SIZE
    run.font.name = "Arial"
    run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height + _TITLE_BAR_GAP, width, _TITLE_BAR_HEIGHT)
    bar.fill.solid()
    bar.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    bar.line.fill.background()
    bar.shadow.inherit = False


def _add_section_divider(prs: Presentation, section: str) -> None:
    """The layout's own default Title position is an off-slide leftover
    (negative top) never actually used by any real slide -- repositioned
    here to match where the reference deck's own real divider slides
    (e.g. "Demographic Trends") placed it. Bold/size/color are likewise
    copied from that real slide's title run: the reference deck's author
    had detached that title into a plain text box with explicit run-level
    formatting rather than relying on placeholder inheritance, so a fresh
    placeholder here renders in the master's plain default (black,
    regular weight) unless the same formatting is applied explicitly."""
    layout = _layout(prs, "3_Large Statement Teal")
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    left, top, width, height = _DIVIDER_TITLE_BOX
    title.left, title.top, title.width, title.height = left, top, width, height
    title.text_frame.text = section
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    run.font.size = _DIVIDER_FONT_SIZE
    run.font.name = "Arial"
    run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1


def _transpose_series(series_by_label: dict) -> dict:
    """{label: {year: value}} -> {year: {label: value}}, the shape
    _category_chart_slide expects -- needed for market-overview volume
    metrics (bar), not for rate metrics (which go straight into
    _multi_line_chart_slide's series_by_label shape unchanged)."""
    years = sorted({y for s in series_by_label.values() for y in s})
    return {y: {label: series_by_label[label].get(y) for label in series_by_label} for y in years}


def _add_heartbeat_slides(prs: Presentation, section: str, rows: list) -> None:
    """rows: costar_heartbeat.parse_properties()'s output. Reuses that
    module's own _decade_class_table (the exact table its one Excel chart
    is built from) rather than recomputing the same decade/class rollup
    here -- one stacked-bar slide, same chart the Excel export produces."""
    decades, by_class = _decade_class_table(rows)
    if not decades:
        return
    categories = {d: {c: by_class[c][d] for c in BROAD_CLASSES} for d in decades}
    _category_chart_slide(prs, section, "Commercial Real Estate Delivered by Decade", categories, "stacked_bar", "count")


def _add_market_overview_slides(prs: Presentation, section: str, markets: list) -> None:
    """markets: [{"name": str, "files": {class: bytes}}], same shape
    costar_market_overview.build_market_overview_workbook takes. One slide
    per (class, metric) that at least one market uploaded -- market name as
    series label, mirroring that module's own per-class-per-metric loop
    exactly (rate metrics as a combined line chart, volume metrics as a
    clustered bar)."""
    for cls in PROPERTY_CLASSES:
        class_data = {}
        for m in markets:
            file_bytes = m["files"].get(cls)
            if file_bytes is None:
                continue
            class_data[m["name"]] = parse_grid(file_bytes, cls)
        if not class_data:
            continue
        for label, _ in METRICS_BY_CLASS[cls]:
            series_by_market = {name: series[label] for name, series in class_data.items() if series.get(label)}
            if not series_by_market:
                continue
            title = f"{CLASS_LABELS[cls]} -- {label}"
            value_format = _market_metric_format(label)
            if _is_rate_metric(label):
                _multi_line_chart_slide(prs, section, title, series_by_market, value_format)
            else:
                _category_chart_slide(prs, section, title, _transpose_series(series_by_market), "bar", value_format)


def _add_smartre_slides(prs: Presentation, section: str, rows: list) -> None:
    """rows: smartre_sales.parse_sales_file() output, already filtered to
    the user's selected subdivisions (same filtering
    build_sales_analysis_workbook itself does). Only the "Overall" cut's
    year x price-bin stacked bar -- see the module docstring's "Known
    simplification (v2, SmartRE)" for why the other 11 cuts and the
    scatter chart aren't reproduced here."""
    dated_rows = [r for r in rows if r["sale_date"]]
    if not dated_rows:
        return
    years = sorted({r["sale_date"].year for r in dated_rows})
    categories = {y: {b: 0 for b in PRICE_BIN_LABELS} for y in years}
    for r in dated_rows:
        categories[r["sale_date"].year][price_bin_label(r["price"])] += 1
    _category_chart_slide(prs, section, "Overall -- Sales by Year & Price Bin", categories, "stacked_bar", "count")


def _build_cre_section(
    prs: Presentation,
    heartbeat_rows: list,
    market_overview_markets: list,
    comparison_costar: list = None,
) -> None:
    """Commercial Real Estate Analysis -- Heartbeat and Market Overview are
    both optional and independent for the subject; the section (and its
    divider) is skipped entirely if the subject uploaded neither AND no
    comparison geography uploaded either. comparison_costar: [(label,
    heartbeat_rows_or_None, market_overview_markets_or_None), ...] -- each
    comparison geography's own optional Heartbeat/Market Overview uploads,
    landing as extra slides under this same section (labeled by geography,
    not a separate divider) rather than in the ACS/BLS-only Comparative
    Analysis section, since CoStar/SmartRE data isn't fetched per-geoid the
    way ACS/BLS dashboards are."""
    comparison_costar = comparison_costar or []
    has_heartbeat = bool(heartbeat_rows)
    has_market = bool(market_overview_markets) and any(m.get("files") for m in market_overview_markets)
    has_comparison = any(
        bool(hb) or (bool(mk) and any(m.get("files") for m in mk)) for _, hb, mk in comparison_costar
    )
    if not has_heartbeat and not has_market and not has_comparison:
        return

    section = "Commercial Real Estate Analysis"
    _add_section_divider(prs, section)
    if has_heartbeat:
        _add_heartbeat_slides(prs, section, heartbeat_rows)
    if has_market:
        _add_market_overview_slides(prs, section, market_overview_markets)

    for label, hb, mk in comparison_costar:
        geo_section = f"{section} -- {label}"
        if hb:
            _add_heartbeat_slides(prs, geo_section, hb)
        if mk and any(m.get("files") for m in mk):
            _add_market_overview_slides(prs, geo_section, mk)


def _comparative_snapshot_categories(geo_categories: list) -> dict:
    """geo_categories: [(label, {year: {series_label: value}}), ...] for
    geographies that have this metric. Returns an ordered
    {f"{year} -- {label}": {series_label: value}} for just the first and
    last year present ACROSS ALL the geographies combined (not per-geo --
    every geo in a single master-deck request already shares the same
    start_year/end_year, so this only diverges from a genuine "first/last"
    when one geo's data has real gaps at the edges), grouped year-major so
    every geography at the same point in time sits next to each other on
    the axis. NOT re-sorted by _add_comparative_snapshot_slide below --
    order here (year-major, geo-minor) is the actual display order."""
    all_years = sorted({y for _, cats in geo_categories for y in cats})
    if not all_years:
        return {}
    snapshot_years = [all_years[0]] if len(all_years) == 1 else [all_years[0], all_years[-1]]
    out = {}
    for y in snapshot_years:
        for label, cats in geo_categories:
            if y in cats:
                out[f"{y} -- {label}"] = cats[y]
    return out


def _add_comparative_snapshot_slide(prs: Presentation, section: str, title: str, chart_type: str, categories_ordered: dict, value_format: str) -> None:
    """categories_ordered: from _comparative_snapshot_categories, already in
    display order -- unlike _category_chart_slide, this does NOT re-sort
    the category keys (they're synthetic "{year} -- {geo}" strings, not
    real years, so a plain sort would scramble the year-major grouping)."""
    cats = list(categories_ordered)
    labels = list(dict.fromkeys(lbl for c in cats for lbl in categories_ordered[c]))
    series = {lbl: [categories_ordered[c].get(lbl) for c in cats] for lbl in labels}
    _add_chart_slide(prs, section, title, chart_type, cats, series, value_format)


def _build_comparative_section(
    prs: Presentation,
    subject_label: str,
    subject_acs: dict,
    subject_bls: dict,
    comparisons: list,
    comparison_acs: list,
    comparison_bls: list,
) -> None:
    """comparisons: [(label, acs_dashboard, bls_dashboard), ...] for each
    comparison geography (subject passed separately, always first).
    comparison_acs/comparison_bls: the chart keys independently selected
    for THIS section (separate state from the single-geography selection
    upstream -- the same subject/comparison dashboards are reused, since
    they're already fetched with every ACS key and whichever BLS sectors
    cover both selections, just filtered differently here).

    "line" metrics (one value per year, e.g. population, median income)
    get one combined multi-series chart, one line per geography, full year
    range -- a real trend is the whole point there. "multi_line"/
    "stacked_bar"/"bar" metrics (multiple values PER YEAR -- sectors or
    category breakdowns) instead collapse to first-year/last-year snapshots
    combined into ONE chart across every geography (confirmed by explicit
    feedback: showing the full year range per geography, one slide each,
    both overwhelmed the deck with slides and buried the actual cross-geo
    comparison the section exists for). "multi_line" data is transposed via
    _transpose_series first so it's in the same {year: {label: value}}
    shape "stacked_bar"/"bar" already use."""
    all_geo = [(subject_label, subject_acs, subject_bls)] + list(comparisons)
    ordered_acs_keys = [k for keys in ACS_SECTION_CHARTS.values() for k in keys if k in comparison_acs]
    if not ordered_acs_keys and not comparison_bls:
        return

    section = "Comparative Analysis"
    _add_section_divider(prs, section)

    for key in ordered_acs_keys:
        if key not in subject_acs:
            continue
        title = acs_chart_title(key)
        chart_type = subject_acs[key]["chart_type"]
        value_format = ACS_CHART_FORMAT.get(key, "count")
        if chart_type == "line":
            series_by_geo = {label: acs[key]["series"] for label, acs, _ in all_geo if key in acs}
            _multi_line_chart_slide(prs, section, title, series_by_geo, value_format)
        else:
            geo_categories = []
            for label, acs, _ in all_geo:
                if key not in acs:
                    continue
                chart = acs[key]
                cats = _transpose_series(chart["series_by_label"]) if chart["chart_type"] == "multi_line" else chart["categories"]
                geo_categories.append((label, cats))
            snapshot = _comparative_snapshot_categories(geo_categories)
            if snapshot:
                snapshot_type = "bar" if chart_type == "multi_line" else chart_type
                _add_comparative_snapshot_slide(prs, section, title, snapshot_type, snapshot, value_format)

    for key in comparison_bls:
        if key not in subject_bls:
            continue
        title = bls_chart_title(key, NAICS_SECTORS)
        chart_type = subject_bls[key]["chart_type"]
        value_format = _bls_chart_format(key)
        if chart_type == "line":
            series_by_geo = {label: bls_d[key]["series"] for label, _, bls_d in all_geo if key in bls_d}
            _multi_line_chart_slide(prs, section, title, series_by_geo, value_format)
        else:
            geo_categories = []
            for label, _, bls_d in all_geo:
                if key not in bls_d:
                    continue
                chart = bls_d[key]
                cats = _transpose_series(chart["series_by_label"]) if chart["chart_type"] == "multi_line" else chart["categories"]
                geo_categories.append((label, cats))
            snapshot = _comparative_snapshot_categories(geo_categories)
            if snapshot:
                snapshot_type = "bar" if chart_type == "multi_line" else chart_type
                _add_comparative_snapshot_slide(prs, section, title, snapshot_type, snapshot, value_format)


def build_master_deck(
    geo_label: str,
    acs_dashboard: dict,
    bls_dashboard: dict,
    selected_acs: list,
    selected_bls: list,
    comparisons: list = None,
    comparison_acs: list = None,
    comparison_bls: list = None,
    heartbeat_rows: list = None,
    market_overview_markets: list = None,
    smartre_rows: list = None,
    comparison_costar: list = None,
    report_title: str = None,
) -> Presentation:
    """geo_label: display name for the title slide, used as a fallback when
    report_title is blank. acs_dashboard/bls_dashboard: already-fetched
    {chart_key: ChartResult} dicts (empty dict if that domain had nothing
    selected). selected_acs/selected_bls: the chart keys the user actually
    opted into -- nothing is included by default, so a chart only gets a
    slide if BOTH the caller fetched it AND the user selected it. One
    section-divider slide per non-empty section; a section with nothing
    selected is skipped entirely. comparisons/comparison_acs/comparison_bls:
    optional Comparative Analysis section, see _build_comparative_section.
    heartbeat_rows/market_overview_markets/comparison_costar: optional
    Commercial Real Estate Analysis section (subject + optionally per-
    comparison-geo), see _build_cre_section -- already-parsed
    (costar_heartbeat.parse_properties output) and raw-bytes-per-class
    (costar_market_overview's own markets shape) respectively, since the
    latter's parsing happens per-class inside _add_market_overview_slides.
    smartre_rows: already-parsed, already-subdivision-filtered
    smartre_sales.parse_sales_file() output (subject only) -- spliced into
    the Housing Analysis section below rather than getting its own divider,
    per the plan's original section mapping. report_title: the user-typed
    name for the report (optional) -- printed on the title slide in place
    of geo_label when given; the caller still uses geo_label for the
    downloaded filename fallback."""
    prs = Presentation(TEMPLATE_PATH)
    _add_title_slide(prs, report_title or geo_label)

    selected_acs_set = set(selected_acs)

    for section, acs_keys in ACS_SECTION_CHARTS.items():
        section_acs = [k for k in acs_keys if k in selected_acs_set and k in acs_dashboard]
        section_bls = list(selected_bls) if section == "Economic Analysis" else []
        section_bls = [k for k in section_bls if k in bls_dashboard]
        section_smartre = smartre_rows if (section == "Housing Analysis" and smartre_rows) else None
        if not section_acs and not section_bls and not section_smartre:
            continue

        _add_section_divider(prs, section)
        for key in section_acs:
            add_dashboard_chart_slide(prs, section, acs_chart_title(key), acs_dashboard[key], ACS_CHART_FORMAT.get(key, "count"))
        for key in section_bls:
            add_dashboard_chart_slide(prs, section, bls_chart_title(key, NAICS_SECTORS), bls_dashboard[key], _bls_chart_format(key))
        if section_smartre:
            _add_smartre_slides(prs, section, section_smartre)

    _build_cre_section(prs, heartbeat_rows or [], market_overview_markets or [], comparison_costar or [])

    if comparisons:
        _build_comparative_section(
            prs, geo_label, acs_dashboard, bls_dashboard, comparisons, comparison_acs or [], comparison_bls or []
        )

    return prs


def pptx_to_bytes(prs: Presentation) -> bytes:
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
